import io
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BRAND = RGBColor(0x1E, 0x3A, 0x8A)
GREY  = RGBColor(0x47, 0x55, 0x69)

def fig_to_png(fig, w=1400, h=800) -> bytes | None:
    try:
        return pio.to_image(fig, format="png", width=w, height=h, scale=2)
    except Exception:
        return None

def fig_to_jpg(fig) -> bytes | None:
    try:
        return pio.to_image(fig, format="jpg", width=1400, height=800, scale=2)
    except Exception:
        return None

def build_pptx(selections: list, title: str = "HR Analytics Report") -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Title slide ─────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BRAND; bg.line.fill.background()
    tf = s.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(11), Inches(1.4)).text_frame
    tf.text = title
    p = tf.paragraphs[0]; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    sf = s.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(11), Inches(0.7)).text_frame
    sf.text = "Generated with the Conglomerate HR Analytics Dashboard"
    sf.paragraphs[0].font.size = Pt(17)
    sf.paragraphs[0].font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)

    # ── One slide per chart ──────────────────────────────────────────────────
    for item in selections:
        img = fig_to_png(item.get("fig"))
        if not img:
            continue
        s = prs.slides.add_slide(blank)
        bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.82))
        bar.fill.solid(); bar.fill.fore_color.rgb = BRAND; bar.line.fill.background()
        ht = s.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.62)).text_frame
        ht.text = item.get("title", "")
        hp = ht.paragraphs[0]; hp.font.size = Pt(22); hp.font.bold = True
        hp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        s.shapes.add_picture(io.BytesIO(img), Inches(0.3), Inches(0.95), height=Inches(5.7))
        if item.get("insight"):
            nb = s.shapes.add_textbox(Inches(0.3), Inches(6.75), Inches(12.5), Inches(0.55)).text_frame
            nb.word_wrap = True
            nb.text = f"💡  {item['insight']}"
            np_ = nb.paragraphs[0]; np_.font.size = Pt(11); np_.font.italic = True
            np_.font.color.rgb = GREY

    out = io.BytesIO(); prs.save(out); out.seek(0)
    return out.getvalue()
